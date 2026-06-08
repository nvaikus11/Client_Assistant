"""Document ingestion for an engagement, organized by material category.

Why this exists: tools need one reliable way to turn the files a consultant drops
into a client's folder into labeled plain text the model can read. Centralizing it
here (per CLAUDE.md) means parsing logic isn't duplicated across tools, and every
tool labels documents the same way.

An *engagement* (one client) is a folder of category subfolders:

    clients/<name>/
        sow-rfp/            SOW / RFP / proposals
        meeting-summaries/  summaries & transcripts (date-stamped filenames)
        exec-updates/       executive updates
        other/              anything else
        outputs/            generated artifacts (NOT an input category)
        <your-folder>/      add more anytime — auto-discovered as a category

`load_engagement()` discovers every subfolder except `outputs/`, parses the
supported files inside, and returns an `Engagement` ready to drop into a prompt.
Adding a new category folder requires no code change. Unsupported files are
reported, not silently skipped — we fail loudly on bad input.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from datetime import date
from enum import Enum
from pathlib import Path


class DocType(str, Enum):
    """How a document should be weighted by downstream prompts."""

    SOW = "SOW"
    RFP = "RFP"
    DECK = "deck"
    NOTES = "notes"
    TRANSCRIPT = "transcript"
    EXEC_UPDATE = "exec update"
    OTHER = "other"


# The reserved output folder is never treated as an input category.
OUTPUTS_DIRNAME = "outputs"

# Category folder -> default DocType for files inside it. None means "infer per
# file" (e.g. sow-rfp, where each file may be a SOW or an RFP). Unknown folders
# default to OTHER, so user-added folders work with no code change.
CATEGORY_DOC_TYPES: dict[str, DocType | None] = {
    "sow-rfp": None,
    "meeting-summaries": DocType.TRANSCRIPT,
    "exec-updates": DocType.EXEC_UPDATE,
    "other": DocType.OTHER,
}

# Preferred display/order for known categories; others follow alphabetically.
_CATEGORY_ORDER = ["sow-rfp", "meeting-summaries", "exec-updates", "other"]

# Folders whose files should be ordered oldest -> newest by filename date.
_DATE_ORDERED = {"meeting-summaries"}

_CATEGORY_LABELS = {
    "sow-rfp": "SOW / RFP",
    "meeting-summaries": "Meeting summaries",
    "exec-updates": "Executive updates",
    "other": "Other",
}


def category_label(folder_name: str) -> str:
    """Human-friendly heading for a category folder name."""
    if folder_name in _CATEGORY_LABELS:
        return _CATEGORY_LABELS[folder_name]
    return folder_name.replace("-", " ").replace("_", " ").strip().title()


@dataclass
class ParsedDoc:
    """Extracted text from one source document, with provenance."""

    filename: str
    doc_type: DocType
    text: str
    category: str = ""
    doc_date: date | None = None  # parsed from filename when present

    def labeled(self) -> str:
        """Render as a labeled block so the model knows each fact's source."""
        when = f" — {self.doc_date.isoformat()}" if self.doc_date else ""
        return f"### {self.filename} [{self.doc_type.value}{when}]\n\n{self.text.strip()}\n"


SUPPORTED_SUFFIXES = {".pdf", ".docx", ".pptx", ".txt", ".md"}
_DATE_IN_NAME = re.compile(r"(\d{4})-(\d{2})-(\d{2})")


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    return "\n".join(p.text for p in document.paragraphs)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        lines.append(f"[Slide {index}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


_READERS = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".txt": _read_text,
    ".md": _read_text,
}


def _infer_doc_type(filename: str) -> DocType:
    """Best-effort guess from the filename; used when a folder has no fixed type."""
    name = filename.lower()
    if "sow" in name or "statement of work" in name:
        return DocType.SOW
    if "rfp" in name or "request for proposal" in name:
        return DocType.RFP
    if name.endswith(".pptx") or "deck" in name or "slides" in name:
        return DocType.DECK
    if "note" in name or "minutes" in name:
        return DocType.NOTES
    return DocType.OTHER


def _date_from_name(filename: str) -> date | None:
    """Pull a YYYY-MM-DD date out of a filename, if present."""
    match = _DATE_IN_NAME.search(filename)
    if not match:
        return None
    try:
        return date(int(match[1]), int(match[2]), int(match[3]))
    except ValueError:
        return None


def parse_document(
    path: Path,
    doc_type: DocType | None = None,
    *,
    category: str = "",
) -> ParsedDoc:
    """Extract text from a single document.

    Raises FileNotFoundError if the file is missing and ValueError if the format
    is unsupported — we never silently proceed on bad input. If `doc_type` is
    None, it is inferred from the filename.
    """
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {path}")

    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise ValueError(
            f"Unsupported file type '{suffix}' for {path.name}. "
            f"Supported: {', '.join(sorted(SUPPORTED_SUFFIXES))}."
        )

    return ParsedDoc(
        filename=path.name,
        doc_type=doc_type or _infer_doc_type(path.name),
        text=_READERS[suffix](path),
        category=category,
        doc_date=_date_from_name(path.name),
    )


def is_supported(path: Path) -> bool:
    """True if the file is a parseable document (and not a placeholder)."""
    return (
        path.is_file()
        and path.suffix.lower() in SUPPORTED_SUFFIXES
        and path.name != ".gitkeep"
    )


def _ordered_categories(client_dir: Path) -> list[str]:
    """Subfolders that count as input categories, in display order."""
    found = [
        p.name
        for p in client_dir.iterdir()
        if p.is_dir()
        and p.name != OUTPUTS_DIRNAME
        and not p.name.startswith((".", "_"))
    ]
    known = [c for c in _CATEGORY_ORDER if c in found]
    extra = sorted(c for c in found if c not in _CATEGORY_ORDER)
    return known + extra


def _parse_category(client_dir: Path, category: str) -> list[ParsedDoc]:
    """Parse every supported file in one category folder, in a sensible order."""
    folder = client_dir / category
    forced = CATEGORY_DOC_TYPES.get(category, DocType.OTHER)
    docs = [
        parse_document(p, doc_type=forced, category=category)
        for p in sorted(folder.iterdir())
        if is_supported(p)
    ]
    if category in _DATE_ORDERED:
        docs.sort(key=lambda d: (d.doc_date or date.max, d.filename))
    return docs


@dataclass
class Engagement:
    """One client's source material, grouped by category, ready for a prompt."""

    client: str
    categories: dict[str, list[ParsedDoc]] = field(default_factory=dict)

    @property
    def all_docs(self) -> list[ParsedDoc]:
        return [doc for docs in self.categories.values() for doc in docs]

    def to_prompt(self) -> str:
        """Render the engagement context as labeled Markdown sections."""
        sections = [f"# ENGAGEMENT SOURCES — CLIENT: {self.client}"]
        for category, docs in self.categories.items():
            heading = category_label(category).upper()
            sections.append(f"\n## {heading}")
            if docs:
                sections += [doc.labeled() for doc in docs]
            else:
                sections.append("_None provided._")
        return "\n\n".join(sections)


def load_engagement(clients_root: Path, client: str) -> Engagement:
    """Load every category folder for one client into an `Engagement`.

    Fails loudly if the client folder is missing, or if it contains no documents
    at all — there would be nothing to ground output on.
    """
    client_dir = clients_root / client
    if not client_dir.exists():
        raise FileNotFoundError(
            f"Client folder not found: {client_dir}. Onboard the client first "
            f"(copy clients/_template, or run scripts/new_client.py)."
        )

    categories = {
        category: _parse_category(client_dir, category)
        for category in _ordered_categories(client_dir)
    }
    engagement = Engagement(client=client, categories=categories)

    if not engagement.all_docs:
        raise ValueError(
            f"No documents found for client '{client}'. Add files to one of the "
            f"category folders under {client_dir} (e.g. sow-rfp/, meeting-summaries/)."
        )
    return engagement


def combine_for_prompt(docs: list[ParsedDoc]) -> str:
    """Concatenate parsed docs into one labeled blob for a user message."""
    return "\n\n".join(doc.labeled() for doc in docs)
